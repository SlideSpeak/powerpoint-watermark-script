import io
import math
from typing import Optional, Tuple

from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def _calculate_ribbon_dimensions(watermark_img: Image.Image, slide_width: int, slide_height: int, 
                               position: str) -> Tuple[int, int]:
    """Calculate watermark dimensions for ribbon positions."""
    original_aspect_ratio = watermark_img.height / watermark_img.width
    width_height_ratio = watermark_img.width / watermark_img.height
    
    # Determine ribbon size based on image aspect ratio
    if width_height_ratio > 2.5:
        # Very wide image - use smaller ribbon height
        size_percentage = 0.20  # 20%
    else:
        # Normal/taller/square image - use larger ribbon height  
        size_percentage = 0.35  # 35%
    
    if position == 'vertical-ribbon':
        # For vertical ribbon, calculate based on width
        ribbon_width = int(slide_width * size_percentage)
        watermark_width = ribbon_width  
        watermark_height = int(ribbon_width * original_aspect_ratio)
    else:
        # For horizontal and diagonal ribbons, calculate based on height
        ribbon_height = int(slide_height * size_percentage)
        watermark_height = ribbon_height
        watermark_width = int(ribbon_height / original_aspect_ratio)
    
    return watermark_width, watermark_height


def _calculate_standard_dimensions(watermark_img: Image.Image, slide_width: int, 
                                 size_percentage: float) -> Tuple[int, int]:
    """Calculate watermark dimensions for standard positions."""
    watermark_width = int(slide_width * size_percentage)
    aspect_ratio = watermark_img.height / watermark_img.width
    watermark_height = int(watermark_width * aspect_ratio)
    return watermark_width, watermark_height


def _get_position_coordinates(position: str, slide_width: int, slide_height: int, 
                            watermark_width: int, watermark_height: int) -> Tuple[int, int]:
    """Calculate the left and top coordinates based on position."""
    positions = {
        'center': (
            (slide_width - watermark_width) // 2,
            (slide_height - watermark_height) // 2
        ),
        'bottom-right': (
            slide_width - watermark_width - Inches(0.5),
            slide_height - watermark_height - Inches(0.5)
        ),
        'bottom-left': (
            Inches(0.5),
            slide_height - watermark_height - Inches(0.5)
        ),
        'top-right': (
            slide_width - watermark_width - Inches(0.5),
            Inches(0.5)
        ),
        'top-left': (
            Inches(0.5),
            Inches(0.5)
        ),
        'diagonal-ribbon': (
            (slide_width - watermark_width) // 2,
            (slide_height - watermark_height) // 2
        ),
        'horizontal-ribbon': (
            0,
            (slide_height - watermark_height) // 2
        ),
        'vertical-ribbon': (
            (slide_width - watermark_width) // 2,
            0
        )
    }
    
    return positions.get(position, positions['center'])


def _process_watermark_image(watermark_path: str, opacity: float) -> Tuple[io.BytesIO, Image.Image]:
    """Process the watermark image with opacity and return as BytesIO stream."""
    watermark_img = Image.open(watermark_path)
    
    # Convert to RGBA if not already
    if watermark_img.mode != 'RGBA':
        watermark_img = watermark_img.convert('RGBA')
    
    # Apply opacity
    alpha = watermark_img.split()[3]
    alpha = alpha.point(lambda p: p * opacity)
    watermark_img.putalpha(alpha)
    
    # Save to bytes stream
    img_stream = io.BytesIO()
    watermark_img.save(img_stream, format='PNG')
    img_stream.seek(0)
    
    return img_stream, watermark_img


def add_watermark_to_pptx(
    pptx_path: str,
    watermark_path: str,
    output_path: Optional[str] = None,
    opacity: float = 0.5,
    position: str = "center",
    size_percentage: float = 0.3,
    on_top: bool = True
) -> str:
    """
    Add a watermark image to all slides in a PowerPoint presentation.
    
    Args:
        pptx_path: Path to the input PowerPoint file
        watermark_path: Path to the watermark image
        output_path: Path for the output file (defaults to input_watermarked.pptx)
        opacity: Opacity of the watermark (0.0 to 1.0, default 0.5)
        position: Position of watermark. Options:
            - Standard: 'center', 'bottom-right', 'bottom-left', 'top-right', 'top-left'
            - Ribbon: 'diagonal-ribbon', 'horizontal-ribbon', 'vertical-ribbon'
        size_percentage: Size of watermark relative to slide width (default 0.3, ignored for ribbon positions)
        on_top: Whether watermark should appear on top of content (True) or underneath (False)
    
    Returns:
        Path to the output file
    """
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Set default output path if not provided
    if output_path is None:
        base_name = pptx_path.rsplit('.', 1)[0]
        output_path = f"{base_name}_watermarked.pptx"
    
    # Process the watermark image
    img_stream, watermark_img = _process_watermark_image(watermark_path, opacity)
    
    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Calculate watermark dimensions based on position type
    if position.endswith('-ribbon'):
        watermark_width, watermark_height = _calculate_ribbon_dimensions(
            watermark_img, slide_width, slide_height, position
        )
    else:
        watermark_width, watermark_height = _calculate_standard_dimensions(
            watermark_img, slide_width, size_percentage
        )
    
    # Get position coordinates
    left, top = _get_position_coordinates(
        position, slide_width, slide_height, watermark_width, watermark_height
    )
    
    # Add watermark to each slide
    for slide in prs.slides:
        # Add the watermark image
        pic = slide.shapes.add_picture(
            img_stream,
            left=left,
            top=top,
            width=watermark_width,
            height=watermark_height
        )
        
        # Apply rotation for diagonal ribbon
        if position == 'diagonal-ribbon':
            angle = math.degrees(math.atan2(slide_height, slide_width))
            pic.rotation = int(angle * 60000)  # PowerPoint uses 60000ths of a degree
        
        # Reset stream position for next slide
        img_stream.seek(0)
        
        # Move watermark to back if requested
        if not on_top:
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
    
    # Save the presentation
    prs.save(output_path)
    return output_path


# Example usage
if __name__ == "__main__":
    # Create a diagonal ribbon watermark
    result = add_watermark_to_pptx(
        pptx_path="presentation.pptx",
        watermark_path="watermark.png",
        output_path="presentation_with_watermark.pptx",
        opacity=0.3,
        position="diagonal-ribbon",
        on_top=True
    )
    print(f"Watermarked presentation saved as: {result}")
    